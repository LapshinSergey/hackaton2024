import hashlib
import os
import nltk
import string
import magic
import PyPDF2
from PyPDF2 import PdfReader
import docx
import pptx
from langdetect import detect
from fastapi import FastAPI, File, UploadFile, WebSocket, HTTPException
from fastapi.responses import HTMLResponse
from fastapi.logger import logger
from elasticsearch import Elasticsearch, exceptions
from sentence_transformers import SentenceTransformer, util
from transformers import MarianMTModel, MarianTokenizer, T5Tokenizer, T5ForConditionalGeneration, AutoModelForSeq2SeqLM, AutoTokenizer
import torch
from datetime import datetime
import nltk
from nltk.tokenize import sent_tokenize
from sklearn.cluster import KMeans
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
from transformers import AutoModelForCausalLM, AutoTokenizer
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import pandas as pd
import io
import torch
from fastapi.responses import StreamingResponse
from typing import List






nltk.download('stopwords')
nltk.download('punkt_tab')

app = FastAPI()

origins = [
    "http://localhost",
    "http://localhost:8080",
    "http://localhost:8001",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)   

# Настройка устройства для использования GPU
device = torch.device("cuda" if torch.cuda.is_available() else "cpu")

# Инициализация моделей
model_sbert = SentenceTransformer('paraphrase-MiniLM-L6-v2').to(device)

# Загрузка основной модели 
model_vicuna = AutoModelForCausalLM.from_pretrained(
    './llama-base_fine-tuned',
    device_map="auto",  # Автоматическое распределение на доступные устройства
    torch_dtype=torch.float16,  # Использование смешанной точности для экономии памяти
    load_in_8bit=False  # Опционально: загрузка модели в 8-битном формате для снижения потребления памяти
)
model_vicuna.to(device)

tokenizer_vicuna = AutoTokenizer.from_pretrained('./llama-base_fine-tuned', use_fast=True)

# Загрузка модели translate
model_marian = MarianMTModel.from_pretrained('Helsinki-NLP/opus-mt-ru-en').to(device)
tokenizer_marian = MarianTokenizer.from_pretrained('Helsinki-NLP/opus-mt-ru-en')

model_marian_ru = MarianMTModel.from_pretrained('Helsinki-NLP/opus-mt-en-ru').to(device)
tokenizer_marian_ru = MarianTokenizer.from_pretrained('Helsinki-NLP/opus-mt-en-ru')

# Загрузка модели t5
model_flan = AutoModelForSeq2SeqLM.from_pretrained('google/flan-t5-base', device_map="auto", torch_dtype="auto").to(device)
tokenizer_flan = AutoTokenizer.from_pretrained('google/flan-t5-base')

# Инициализация Elasticsearch с указанием хоста
# URL Elasticsearch и данные для авторизации
es_url = 'https://localhost:9200'
username = 'elastic'
password = 'Aetoa1Tahyor'

es = Elasticsearch(
    es_url,
    basic_auth=(username, password),  # Передаем имя пользователя и пароль
    verify_certs=False  # Отключить проверку SSL (если необходимо)
)

# Модели для создания индексов
uploaded_files_mapping = {
    "mappings": {
        "properties": {
            "file_id": {"type": "keyword"},
            "file_name": {"type": "text"},
            "checksum": {"type": "keyword"},
            "uploaded_at": {"type": "date"},
            "file_size": {"type": "long"}
        }
    }
}

keywords_mapping = {
    "mappings": {
        "properties": {
            "keyword": {"type": "text"},
            "file": {"type": "text"},
            "pages": {"type": "keyword"},
            "score": {"type": "float"}
        }
    }
}

learning_mapping = {
    "mappings": {
        "properties": {
            "question": {"type": "text"},
            "answer": {"type": "text"},
            "score": {"type": "float"}
        }
    }
}

learning_score_mapping = {
    "mappings": {
        "properties": {
            "question": {"type": "text"},
            "answer": {"type": "text"},
            "score": {"type": "float"}
        }
    }
}

# Функция для создания индекса, если он не существует
def create_index_if_not_exists(index_name: str, mapping: dict):
    if not es.indices.exists(index=index_name):
        print(f"Индекс {index_name} не существует, создаём его...")
        es.indices.create(index=index_name, body=mapping)
        print(f"Индекс {index_name} успешно создан.")

# Функция для создания всех необходимых индексов
def create_required_indexes():
    create_index_if_not_exists("uploaded_files", uploaded_files_mapping)
    create_index_if_not_exists("keywords", keywords_mapping)
    create_index_if_not_exists("learning", learning_mapping)
    create_index_if_not_exists("learning-score", learning_score_mapping)

# Вызываем функцию для создания индексов при старте приложения
@app.on_event("startup")
async def startup():
    create_required_indexes()

# Настройка логирования
import logging
logger.setLevel(logging.INFO)
handler = logging.StreamHandler()
handler.setFormatter(logging.Formatter('%(levelname)s:%(name)s:%(message)s'))
logger.addHandler(handler)

# Допустимые расширения и MIME-типы файлов
ALLOWED_EXTENSIONS = ["pptx", "doc", "docx", "pdf", "txt"]
ALLOWED_MIME_TYPES = [
    "application/pdf", "text/plain",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
]

# Функция для извлечения текста из PDF
def extract_text_from_pdf(file):
    text = ""
    reader = PyPDF2.PdfReader(file)
    for page in reader.pages:
        text += page.extract_text()
    return text

# Функция для извлечения текста из Word
def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

# Функция для извлечения текста из PPTX
def extract_text_from_pptx(file):
    prs = pptx.Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

# Функция очистки текста
def clean_text(text):
    return text.strip()

# Перевод
def translate_text(text, model=model_marian, tokenizer=tokenizer_marian, device=device, language='english'):
    # lang_codes = {
    #     'ru': 'ru_RU',
    #     'en': 'en_XX'
    # }
    # target_lang = lang_codes[language]

    sentences = sent_tokenize(text, language)
    
    translated_sentences = []
    
    for sentence in sentences:
        # Токенизация предложения
        inputs = tokenizer(sentence, return_tensors="pt", padding=True, truncation=True).to(device)
        
        # Генерация перевода
        with torch.no_grad():
            translated = model.generate(**inputs)
        
        # Декодирование перевода
        translated_sentence = tokenizer.decode(translated[0], skip_special_tokens=True)
        
        translated_sentences.append(translated_sentence)
    
    # Собираем переведённые предложения обратно в текст
    translated_text = " ".join(translated_sentences)
    
    return translated_text

# Функция для сегментации текста с учётом номеров страниц
def sbert_segmentation_with_page_numbers(text, similarity_threshold=0.7):
    # Разделяем текст на страницы по символу "*_*"
    pages = text.split('*_*')
    
    segments = []
    current_segment = []
    current_page = None
    embeddings = []
    
    # Пройдем по каждой странице текста
    for ipage, page_text in enumerate(pages):
        # Разбиваем страницу на предложения или фрагменты (зависит от вашей задачи)
        sentences = page_text.split('. ')  # Разбиваем по предложениям
        
        for sentence in sentences:
            # Получаем эмбеддинг предложения
            embedding = model_sbert.encode(sentence)
            embeddings.append(embedding)
            
            if current_segment:
                # Вычисляем схожесть с последним фрагментом
                last_embedding = embeddings[-2]  # Последний добавленный фрагмент
                similarity = cosine_similarity([embedding], [last_embedding])[0][0]
                
                if similarity < similarity_threshold:
                    # Если схожесть ниже порога, начинаем новый блок
                    segments.append({"text": " ".join(current_segment), "page": ipage+1})
                    current_segment = []  # Очищаем текущий сегмент
            
            # Добавляем текущее предложение в сегмент
            current_segment.append(sentence)
        
        # Добавляем последний сегмент на странице
        if current_segment:
            segments.append({"text": " ".join(current_segment), "page": current_page})
        
        # Очищаем текущий сегмент для следующей страницы
        current_segment = []
    
    return segments

def extract_text_and_pages_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    # page_numbers = []
    
    for i, page in enumerate(reader.pages):
        text += page.extract_text() + "*_*"
        # page_numbers.append(i + 1)  # Номера страниц начинаются с 1
    
    return text

def extract_keywords(text, top_n=10):
    # Токенизация текста на предложения и слова
    sentences = nltk.sent_tokenize(text)
    words = nltk.word_tokenize(text)

    # Фильтрация слов (удаление стоп-слов и знаков препинания)
    words = [word.lower() for word in words if word.isalnum()]

    # Получение эмбеддингов для предложений
    sentence_embeddings = model_sbert.encode(sentences, convert_to_tensor=True)

    # TF-IDF для определения весов слов в контексте
    vectorizer = TfidfVectorizer()
    tfidf_matrix = vectorizer.fit_transform(sentences)
    tfidf_scores = {word: tfidf_matrix.getcol(idx).sum() for word, idx in vectorizer.vocabulary_.items()}

    # Получение эмбеддингов слов
    word_embeddings = model_sbert.encode(list(tfidf_scores.keys()), convert_to_tensor=True)

    # Вычисление косинусного сходства между эмбеддингами предложений и эмбеддингами слов
    average_sentence_embedding = sentence_embeddings.mean(dim=0)
    cosine_scores = util.cos_sim(average_sentence_embedding, word_embeddings)[0]

    # Сортировка слов по косинусному сходству и TF-IDF весу
    top_indices = cosine_scores.argsort(descending=True)[:top_n * 2]
    candidate_keywords = [
        (list(tfidf_scores.keys())[idx], cosine_scores[idx].item() * tfidf_scores[list(tfidf_scores.keys())[idx]]) for
        idx in top_indices]

    # Сортировка по значимости и отбор топ-N ключевых слов
    candidate_keywords.sort(key=lambda x: x[1], reverse=True)
    keywords = [word for word, score in candidate_keywords[:top_n]]

    return keywords 

# Функция для индексации ключевых слов в Elasticsearch
def index_keywords_to_elasticsearch(keywords, file_id, file_name):
    for keyword in keywords:
        # Индексация каждого ключевого слова
        doc = {
            "keyword": keyword,
            "file_id": file_id,
            "file_name": file_name,
            "score": 0.0,  # Можно вычислить вес или важность ключевого слова
            "timestamp": datetime.utcnow().isoformat()
        }
        es.index(index="keywords", body=doc)

# API-метод upload
@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    if file.filename.split(".")[-1] not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Недопустимый формат файла.")
    
    content = await file.read()
    mime = magic.from_buffer(content[:1024], mime=True)
    if mime not in ALLOWED_MIME_TYPES:
        raise HTTPException(status_code=400, detail="Недопустимый MIME-тип файла.")
    
    checksum = hashlib.md5(content).hexdigest()

    # Проверка на дубликат в Elasticsearch
    query = {
        "query": {
            "match": {
                "checksum": checksum
            }
        }
    }
    result = es.search(index="uploaded_files", body=query)
    if result['hits']['total']['value'] > 0:
        raise HTTPException(status_code=400, detail="Файл уже загружен.")

    uploaded_at = datetime.utcnow().isoformat()

    # Сохранение информации о файле
    file_info = {
        "file_id": checksum,
        "file_name": file.filename,
        "checksum": checksum,
        "uploaded_at": uploaded_at,
        "file_size": len(content)
    }
    es.index(index="uploaded_files", body=file_info)

    # Извлечение текста и номеров страниц из файла
    if mime == "application/pdf":
        text = extract_text_and_pages_from_pdf(file.file)
    elif mime in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword"
    ]:
        text = extract_text_from_docx(file.file)
        page_numbers = []  # Для DOCX страницы не поддерживаются
    elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = extract_text_from_pptx(file.file)
        page_numbers = []  # Для PPTX страницы не поддерживаются
    elif mime == "text/plain":
        text = content.decode('utf-8')
        page_numbers = []  # Для TXT страниц не поддерживаются
    else:
        raise HTTPException(status_code=400, detail="Не удалось обработать файл.")

    text = clean_text(text)

    # Определение языка и перевод
    language = detect(text)
    if language == 'ru':
        text = translate_text(text) 

    # Извлечение ключевых слов и индексация
    keywords = extract_keywords(text)
    for keyword in keywords:
        es.index(index="keywords", body={"keyword": keyword})
        
    # Сегментация текста с использованием SBERT
    segments = sbert_segmentation_with_page_numbers(text)

    # Индексация сегментов в Elasticsearch
    for i, segment in enumerate(segments):
        embedding = model_sbert.encode(segment['text']).tolist()
        doc = {
            "file_id": checksum,
            "file_name": file.filename,
            "page_number": segment['page'],
            "segment_id": i,
            "text": segment['text'],
            "embedding": embedding
        }
        es.index(index="text_segments", body=doc)

    return {"filename": file.filename, "checksum": checksum}

# WebSocket метод ask
@app.websocket("/ask")
async def ask_endpoint(websocket: WebSocket):
    await websocket.accept()
    while True:
        question = await websocket.receive_text()
        # await websocket.send_text("Думаю")

        # Определение языка вопроса
        language = detect(question)
        
        corrected_question = question

        if language == 'ru':
            corrected_question = translate_text(corrected_question, model=model_marian, tokenizer=tokenizer_marian, language='english')

        # Поиск релевантных блоков текста
        question_embedding = model_sbert.encode(corrected_question).tolist()
        script_query = {
            "script_score": {
                "query": {"match_all": {}},
                "script": {
                    "source": "cosineSimilarity(params.query_vector, 'embedding') + 1.0",
                    "params": {"query_vector": question_embedding}
                }
            }
        }
        response = es.search(
            index="text_segments",
            body={
                "size": 10,  # Ограничиваем количество найденных сегментов
                "query": script_query,
                "_source": ["file_name", "text", "page_number"]
            }
        )
        relevant_blocks = []
        context = ""
        
        min_score = 1.6

        if response['hits']['hits'] and response['hits']['hits'][0]["_score"] > min_score:
            context = response['hits']['hits'][0]["_source"]["text"]
            for hit in response['hits']['hits']:
                block = {
                    "file": hit["_source"]["file_name"],
                    "pages": [hit["_source"]["page_number"]],
                    "score": hit["_score"],
                    "squeeze": hit["_source"]["text"]
                }
                if hit["_score"] > min_score:
                    relevant_blocks.append(block)
            input_text_flan = f"Question: {corrected_question} Context (if relevant): {context} Answer: "
            # input_text = f"Question: {corrected_question} Context (if relevant): {context} Answer: "
            input_text = context + "User: " + corrected_question + "\nAssistant:"
        else:
            input_text_flan = (
                f"Question: {corrected_question} "
                "Answer the question based on your general knowledge without any specific context. "
                "If relevant, provide as detailed and informative an answer as possible. Answer: "
            )
            input_text = "User: " + corrected_question + "\nAssistant:"
        
        # Генерация ответа с использованием модели FLAN
        input_ids = tokenizer_flan(input_text_flan, return_tensors='pt').input_ids.to(device)
        outputs = model_flan.generate(
            input_ids,
            temperature=0.7,
            max_length=256,
            max_new_tokens = 50,
            top_p=0.9,
            top_k=40,
            do_sample=True,
            use_cache=False,
        )
        answer_flan = tokenizer_flan.decode(outputs[0], skip_special_tokens=True)

        # Генерация ответа с использованием модели VICUNA
        input_ids = tokenizer_vicuna(input_text, return_tensors='pt').input_ids.to(device)
        outputs = model_vicuna.generate(
            input_ids,
            temperature=0.7,
            max_length=256,
            max_new_tokens = 50,
            top_p=0.9,
            top_k=40,
            do_sample=True,
            use_cache=False,
            # no_repeat_ngram_size=2
        )
        answer = tokenizer_vicuna.decode(outputs[0], skip_special_tokens=True)

        index = answer.find("Answer:")
        if index != -1:
            answer = answer.split("Answer:")[-1].strip()
        else:
            answer = answer.split("Assistant:")[-1].strip()

        
        # Перевод ответа на русский при необходимости
        if language == 'ru':
            answer_flan = translate_text(answer_flan, model=model_marian_ru, tokenizer=tokenizer_marian_ru, language='russian')
            answer = translate_text(answer, model=model_marian_ru, tokenizer=tokenizer_marian_ru, language='russian')

        # Формирование ответа
        response = {
            "answer_flan": answer_flan,
            "answer": answer,
            "relevance": relevant_blocks
        }

        await websocket.send_json(response)

@app.post("/process_questions")
async def process_questions(file: UploadFile = File(...)):
    # Read the uploaded Excel file into a pandas DataFrame
    df = pd.read_excel(file.file)

    # Iterate over each question in the DataFrame
    for index, row in df.iterrows():
        question = row['question']

        # Detect the language of the question
        language = detect(question)

        corrected_question = question

        if language == 'ru':
            corrected_question = translate_text(
                corrected_question,
                model=model_marian,
                tokenizer=tokenizer_marian,
                language='english'
            )

        # Generate question embedding
        question_embedding = model_sbert.encode(corrected_question).tolist()

        # Build Elasticsearch query
        script_query = {
            "script_score": {
                "query": {"match_all": {}},
                "script": {
                    "source": "cosineSimilarity(params.query_vector, 'embedding') + 1.0",
                    "params": {"query_vector": question_embedding}
                }
            }
        }

        response = es.search(
            index="text_segments",
            body={
                "size": 1,  # Limit to only the top result
                "query": script_query,
                "_source": ["file_name", "text", "page_number"]
            }
        )

        context = ""
        min_score = 1.6

        if response['hits']['hits'] and response['hits']['hits'][0]["_score"] > min_score:
            hit = response['hits']['hits'][0]
            context = hit["_source"]["text"]
            # Remove file extension from filename
            filename = os.path.splitext(hit["_source"]["file_name"])[0]
            slide_number = hit["_source"]["page_number"]
            # input_text = f"Question: {corrected_question} Context (if relevant): {context} Answer: "
            input_text = context + "User: " + corrected_question + "\nAssistant:"
        else:
            filename = ''
            slide_number = ''
            input_text = "User: " + corrected_question + "\nAssistant:"


        # Generate answer using the language model
        input_ids = tokenizer_vicuna(input_text, return_tensors='pt').input_ids.to(device)
        outputs = model_vicuna.generate(
            input_ids,
            temperature=0.7,
            max_length=256,
            max_new_tokens=50,
            top_p=0.9,
            top_k=40,
            do_sample=True,
            use_cache=False,
        )
        answer = tokenizer_vicuna.decode(outputs[0], skip_special_tokens=True)
        
        index = answer.find("Answer:")
        if index != -1:
            answer = answer.split("Answer:")[-1].strip()
        else:
            answer = answer.split("Assistant:")[-1].strip()

        # Translate the answer back to Russian if necessary
        if language == 'ru':
            answer = translate_text(
                answer,
                model=model_marian_ru,
                tokenizer=tokenizer_marian_ru,
                language='russian'
            )

        # Fill in the 'answer' column
        df.at[index, 'answer'] = answer

        # Update 'filename' and 'slide_number' columns
        df.at[index, 'filename'] = filename
        df.at[index, 'slide_number'] = slide_number

    # Save the DataFrame to an Excel file in memory
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False)
    output.seek(0)

    # Return the Excel file as a streaming response
    return StreamingResponse(
        output,
        media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        headers={"Content-Disposition": "attachment; filename=processed_questions.xlsx"}
    )

# WebSocket метод prompt
@app.websocket("/prompt")
async def prompt_endpoint(websocket: WebSocket):
    await websocket.accept()
    while True:
        prefix = await websocket.receive_text()

        # Поиск ключевых слов по префиксу
        response = es.search(
            index="keywords",
            body={
                "size": 10,
                "query": {
                    "prefix": {
                        "keyword": {
                            "value": prefix
                        }
                    }
                }
            }
        )
        suggestions = [hit["_source"]["keyword"] for hit in response["hits"]["hits"]]

        await websocket.send_json(suggestions)

# Модель данных для принятия JSON
class LearningScore(BaseModel):
    question: str
    answer: str
    score: int

# Метод vote
@app.post("/vote")
async def vote_endpoint(data: LearningScore):
    try:
        # Проверяем, существует ли индекс, если нет, создаем его
        if not es.indices.exists(index="learning-score"):
            es.indices.create(index="learning-score")
        
        # Добавляем запись в индекс
        response = es.index(index="learning-score", document=data.dict())
        return {"status": "success", "data": response}

    except exceptions.ConnectionError:
        raise HTTPException(status_code=500, detail="Connection to Elasticsearch failed")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
